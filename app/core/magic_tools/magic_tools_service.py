import httpx
import uuid
import io
import csv
from urllib.parse import urlparse
from datetime import datetime
from typing import Dict, Any, Optional
from app.core.ocr_service import ocr_service
from app.core.database.s3_service import s3_service

# Try to import optional libraries for Office documents
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("Warning: python-docx not available. .docx files will not be supported.")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not available. .pptx files will not be supported.")

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False
    print("Warning: openpyxl not available. Excel files (.xlsx, .xls) will not be supported.")


class MagicToolsService:
    """Service for handling magic tools operations like file content extraction"""
    
    def __init__(self):
        """Initialize the magic tools service"""
        self.ocr_service = ocr_service
        self.s3_service = s3_service
    
    async def extract_content_from_url(self, file_url: str) -> Dict[str, Any]:
        """
        Extract text content from a file URL.
        
        This method:
        1. Downloads the file from the URL
        2. Extracts text content based on file type (PDF, image, text)
        3. Saves the file to S3 in sanad-data-source bucket
        4. Returns extracted content and S3 information
        
        Args:
            file_url: URL of the file to extract content from
        
        Returns:
            Dictionary with:
            - success: bool
            - content: str (extracted text)
            - filename: str
            - content_type: str
            - s3_bucket: str (sanad-data-source)
            - s3_key: str
            - s3_url: str (presigned or CDN URL)
            - content_length: int
        
        Raises:
            Exception: If download, extraction, or S3 upload fails
        """
        # Download file from URL
        print(f"Downloading file from URL: {file_url}")
        
        async with httpx.AsyncClient(timeout=60.0) as client:
            try:
                response = await client.get(file_url)
                response.raise_for_status()
                file_content = response.content
                content_type = response.headers.get("content-type", "application/octet-stream")
            except httpx.HTTPError as e:
                raise Exception(f"Failed to download file from URL: {str(e)}")
        
        if not file_content:
            raise Exception("Downloaded file is empty")
        
        # Extract filename from URL
        parsed_url = urlparse(file_url)
        filename = parsed_url.path.split('/')[-1] or f"extracted_file_{uuid.uuid4().hex[:8]}"
        
        # Get file extension
        file_ext = ""
        if '.' in filename:
            file_ext = filename.lower().split('.')[-1]
        
        # If filename doesn't have extension, try to infer from content-type
        if not file_ext:
            if 'pdf' in content_type.lower():
                filename += '.pdf'
                file_ext = 'pdf'
            elif 'image' in content_type.lower():
                if 'jpeg' in content_type.lower() or 'jpg' in content_type.lower():
                    filename += '.jpg'
                    file_ext = 'jpg'
                elif 'png' in content_type.lower():
                    filename += '.png'
                    file_ext = 'png'
                else:
                    filename += '.img'
                    file_ext = 'img'
            elif 'text' in content_type.lower():
                filename += '.txt'
                file_ext = 'txt'
            elif 'word' in content_type.lower() or 'document' in content_type.lower():
                filename += '.docx'
                file_ext = 'docx'
            elif 'excel' in content_type.lower() or 'spreadsheet' in content_type.lower():
                filename += '.xlsx'
                file_ext = 'xlsx'
            elif 'presentation' in content_type.lower() or 'powerpoint' in content_type.lower():
                filename += '.pptx'
                file_ext = 'pptx'
        
        # Extract text content based on file extension
        extracted_text = ""
        
        try:
            # Supported file types: .pdf, .txt, .docx, .md, .pptx, .tex, .ts, .csv, .xlsx, .xls
            if file_ext == 'pdf':
                # Use OCR service for PDF
                pages_results = self.ocr_service.extract_text_from_pdf_bytes(file_content)
                extracted_text = "\n\n".join([
                    page.get("raw_text", "") 
                    for page in pages_results 
                    if page.get("raw_text")
                ])
            elif file_ext in ['txt', 'md', 'tex', 'ts']:
                # Text-based files: direct text extraction
                extracted_text = file_content.decode("utf-8", errors="ignore")
            elif file_ext == 'docx':
                # Word documents
                if not DOCX_AVAILABLE:
                    raise Exception("python-docx library is required for .docx files. Please install: pip install python-docx")
                doc = DocxDocument(io.BytesIO(file_content))
                paragraphs = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        paragraphs.append(para.text)
                extracted_text = "\n".join(paragraphs)
                
                # Also extract text from tables
                table_texts = []
                for table in doc.tables:
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            table_texts.append(" | ".join(row_text))
                if table_texts:
                    extracted_text += "\n\n--- Tables ---\n" + "\n".join(table_texts)
            elif file_ext == 'pptx':
                # PowerPoint presentations
                if not PPTX_AVAILABLE:
                    raise Exception("python-pptx library is required for .pptx files. Please install: pip install python-pptx")
                prs = Presentation(io.BytesIO(file_content))
                slides_text = []
                for i, slide in enumerate(prs.slides, 1):
                    slide_content = []
                    slide_content.append(f"--- Slide {i} ---")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_content.append(shape.text.strip())
                    if len(slide_content) > 1:  # More than just the header
                        slides_text.append("\n".join(slide_content))
                extracted_text = "\n\n".join(slides_text)
            elif file_ext == 'csv':
                # CSV files
                try:
                    csv_text = file_content.decode("utf-8", errors="ignore")
                    # Parse CSV and format as readable text
                    csv_reader = csv.reader(io.StringIO(csv_text))
                    rows = []
                    for row in csv_reader:
                        rows.append(" | ".join(row))
                    extracted_text = "\n".join(rows)
                except Exception as csv_error:
                    # Fallback: just decode as text
                    extracted_text = file_content.decode("utf-8", errors="ignore")
            elif file_ext in ['xlsx', 'xls']:
                # Excel files
                if not EXCEL_AVAILABLE:
                    raise Exception("openpyxl library is required for Excel files. Please install: pip install openpyxl")
                workbook = openpyxl.load_workbook(io.BytesIO(file_content), data_only=True)
                sheets_text = []
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    sheet_content = [f"--- Sheet: {sheet_name} ---"]
                    for row in sheet.iter_rows(values_only=True):
                        row_values = [str(cell) if cell is not None else "" for cell in row]
                        if any(row_values):  # Only add non-empty rows
                            sheet_content.append(" | ".join(row_values))
                    if len(sheet_content) > 1:  # More than just the header
                        sheets_text.append("\n".join(sheet_content))
                extracted_text = "\n\n".join(sheets_text)
            elif content_type.startswith("text/"):
                # Generic text content type
                extracted_text = file_content.decode("utf-8", errors="ignore")
            elif content_type.startswith("image/"):
                # Images - use OCR service
                result = self.ocr_service.extract_text_from_image_bytes(file_content)
                extracted_text = result.get("raw_text", "")
            else:
                # Try to decode as text for unknown types
                try:
                    extracted_text = file_content.decode("utf-8", errors="ignore")
                except:
                    raise Exception(
                        f"Unsupported file type: {file_ext or content_type}. "
                        f"Supported types: .pdf, .txt, .docx, .md, .pptx, .tex, .ts, .csv, .xlsx, .xls"
                    )
        except Exception as e:
            raise Exception(f"Error extracting text content: {str(e)}")
        
        if not extracted_text or not extracted_text.strip():
            raise Exception("No text content could be extracted from the file")
        
        # Save file to S3 in sanad-data-source bucket
        s3_key = None
        s3_url = None
        if self.s3_service and self.s3_service.api_configured:
            try:
                bucket_name = "sanad-data-source"
                # Use a folder structure: extracted/{timestamp}/{filename}
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                s3_key = f"extracted/{timestamp}/{filename}"
                
                # Upload to S3
                self.s3_service.s3_client.put_object(
                    Bucket=bucket_name,
                    Key=s3_key,
                    Body=file_content,
                    ContentType=content_type
                )
                print(f"File saved to S3: s3://{bucket_name}/{s3_key}")
                
                # Generate presigned URL or CDN URL if available
                try:
                    if self.s3_service.cdn_url:
                        s3_url = self.s3_service.get_cdn_url(s3_key)
                    else:
                        s3_url = self.s3_service.get_presigned_url(
                            bucket_name=bucket_name,
                            key=s3_key,
                            expiration=3600 * 24 * 7  # 7 days
                        )
                except Exception as url_error:
                    print(f"Warning: Could not generate URL for S3 file: {url_error}")
                    s3_url = f"s3://{bucket_name}/{s3_key}"
            except Exception as s3_error:
                print(f"Warning: Failed to save file to S3: {s3_error}")
                # Continue even if S3 upload fails
        
        return {
            "success": True,
            "content": extracted_text,
            "filename": filename,
            "content_type": content_type,
            "s3_bucket": "sanad-data-source" if s3_key else None,
            "s3_key": s3_key,
            "s3_url": s3_url,
            "content_length": len(extracted_text)
        }


# Create a singleton instance
magic_tools_service = MagicToolsService()

